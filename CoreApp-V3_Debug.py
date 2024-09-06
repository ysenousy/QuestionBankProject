import os
import gc
import re
import torch
import random
import numpy as np
from pptx import Presentation
from docx import Document
from sense2vec import Sense2Vec
from sentence_transformers import SentenceTransformer
import yake
import nltk
import string
from nltk.corpus import stopwords
from sklearn.feature_extraction.text import TfidfVectorizer
from concurrent.futures import ThreadPoolExecutor
from sklearn.metrics.pairwise import cosine_similarity
import spacy

nltk.download('punkt')
nltk.download('wordnet')
nltk.download('stopwords')
nltk.download('omw-1.4')

# Load SpaCy model
spacy_model = spacy.load('en_core_web_sm')

# Initialize models
def initialize_models():
    device = torch.device("cuda" if torch.cuda.is_available() else "cpu")
    print(f"Using device: {device}")

    # Load Sense2Vec model
    s2v = Sense2Vec().from_disk(r'C:\Python\Lib\site-packages\sense2vec\tests\data')

    # Load Sentence Transformer model
    sentence_transformer_model = SentenceTransformer('sentence-transformers/msmarco-distilbert-base-v2')

    return s2v, sentence_transformer_model, device

# Create necessary folders
def create_folders(base_directory):
    os.makedirs(os.path.join(base_directory, "Quizzes"), exist_ok=True)
    os.makedirs(os.path.join(base_directory, "Question Bank"), exist_ok=True)
    os.makedirs(os.path.join(base_directory, "Weekly Assignments"), exist_ok=True)

# Read all text from a PowerPoint file
def read_pptx(file_path):
    presentation = Presentation(file_path)
    all_text = []

    for slide in presentation.slides:
        slide_text = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                slide_text.append(paragraph.text)
        slide_content = "\n".join(slide_text)
        all_text.append(slide_content)

    return all_text

# Clean input text
def clean_text(text):
    text = re.sub(r'[^\x09\x0A\x0D\x20-\x7F]', '', text)
    text = re.sub(r'\d{1,2}/\d{1,2}/\d{2,4}', '', text)
    text = re.sub(r'\s+', ' ', text)
    return text.strip()

# TF-IDF keyword extraction
def get_tfidf_keywords(content, top_n=20):
    content_cleaned = [re.sub(r'\s+', ' ', content.lower())]
    vectorizer = TfidfVectorizer(stop_words='english')
    tfidf_matrix = vectorizer.fit_transform(content_cleaned)
    feature_names = vectorizer.get_feature_names_out()
    tfidf_scores = tfidf_matrix.toarray().flatten()
    top_indices = tfidf_scores.argsort()[-top_n:][::-1]
    tfidf_keywords = [(feature_names[i], tfidf_scores[i]) for i in top_indices]
    return [keyword for keyword, score in tfidf_keywords]

# YAKE keyword extraction
def get_yake_keywords(content, top_n=10):
    kw_extractor = yake.KeywordExtractor(lan="en", n=2, dedupLim=0.9, top=top_n)
    yake_keywords = kw_extractor.extract_keywords(content)
    return [kw for kw, score in yake_keywords]

# Combine TF-IDF and YAKE keywords and filter by POS using SpaCy
def get_keywords(content):
    try:
        # Extract TF-IDF and YAKE keywords
        tfidf_keywords = get_tfidf_keywords(content, top_n=20)
        yake_keywords = get_yake_keywords(content, top_n=10)

        # Combine keywords and prioritize multi-word phrases
        combined_keywords = list(set(tfidf_keywords + yake_keywords))
        multi_word_keywords = [k for k in combined_keywords if len(k.split()) > 1]
        single_word_keywords = [k for k in combined_keywords if len(k.split()) == 1]

        # Filter by part of speech using SpaCy
        final_keywords = filter_by_pos(multi_word_keywords + single_word_keywords)

        # Debug: Print extracted keywords
        print(f"Extracted keywords: {final_keywords}")
        
        return final_keywords[:15]

    except Exception as e:
        print(f"Error in keyword extraction: {e}")
        return []

# Function to filter keywords by part-of-speech using SpaCy
def filter_by_pos(keywords):
    filtered_keywords = []
    for keyword in keywords:
        doc = spacy_model(keyword)
        if any(token.pos_ in {'NOUN', 'PROPN', 'ADJ'} for token in doc):
            filtered_keywords.append(keyword)
    return filtered_keywords

# Generate distractors using Sense2Vec
def sense2vec_get_words(word, s2v, top_n=10, sentence=None):
    word = word.lower()
    word_with_pos = f"{word}|NOUN"
    if word_with_pos in s2v:
        most_similar = s2v.most_similar(word_with_pos, n=top_n)
        return [w[0].split("|")[0] for w in most_similar]
    else:
        return []

# Enhanced distractor generation
def get_distractors(word, sentence, s2v, sentence_transformer_model, top_n=40, lambda_val=0.2):
    distractors = sense2vec_get_words(word, s2v, top_n, sentence)
    doc = spacy_model(sentence)
    word_pos = next((token.pos_ for token in doc if token.text.lower() == word), None)
    word_ner = next((ent.label_ for ent in doc.ents if word in ent.text.lower()), None)
    filtered_distractors = [distractor for distractor in distractors if spacy_model(distractor)[0].pos_ == word_pos]
    
    # Debug: Print generated distractors
    print(f"Distractors for {word}: {filtered_distractors}")

    return filtered_distractors[:3] if filtered_distractors else distractors[:3]

# Process a single slide (parallel processing)
def process_slide(slide_text, s2v, sentence_transformer_model, device, question_count, max_questions=20):
    cleaned_text = clean_text(slide_text)
    imp_keywords = get_keywords(cleaned_text)

    slide_questions = []
    for answer in imp_keywords:
        if question_count[0] >= max_questions:  # Stop when max_questions are reached
            break
        question = generate_question_templates(answer)

        # Debug: Print generated question
        print(f"Generated question: {question}")

        distractors = get_distractors(answer.capitalize(), question, s2v, sentence_transformer_model)
        if len(distractors) < 3:
            continue

        options = distractors[:3] + [answer]
        random.shuffle(options)
        correct_option = options.index(answer)

        output = f"Question: {question}\n"
        for i, opt in enumerate(options):
            output += f"{chr(65 + i)}) {opt}\n"
        output += f"Answer: {chr(65 + correct_option)}\n\n"
        slide_questions.append(output)
        question_count[0] += 1  # Increment the global question counter

    return slide_questions

# Generate all questions for a PPTX
def generate_questions_for_pptx(pptx_path, output_directory, s2v, sentence_transformer_model, device, max_questions=20):
    all_questions = []
    question_count = [0]  # Use list to allow modification inside process_slide
    slide_texts = read_pptx(pptx_path)

    with ThreadPoolExecutor(max_workers=4) as executor:
        results = executor.map(lambda slide_text: process_slide(slide_text, s2v, sentence_transformer_model, device, question_count, max_questions), slide_texts)
        for slide_questions in results:
            all_questions.extend(slide_questions)
            if question_count[0] >= max_questions:
                break

    if all_questions:
        full_text = "\n\n".join(all_questions)
        word_filename = os.path.basename(pptx_path).replace('.pptx', '.docx')
        word_path = os.path.join(output_directory, "Quizzes", word_filename)
        create_word_document(full_text, word_path)
        print(f"Processed {pptx_path} -> {word_path}")
    else:
        print(f"No questions generated for {pptx_path}")

# Utility functions for formatting text and saving documents
def create_word_document(questions, output_path):
    doc = Document()
    doc.add_heading('Generated Questions', 0)
    doc.add_paragraph(questions)
    doc.save(output_path)

# Maximal Marginal Relevance function
def mmr(doc_embedding, word_embeddings, words, top_n, lambda_param=0.5):
    word_doc_similarity = cosine_similarity(word_embeddings, doc_embedding)
    word_similarity = cosine_similarity(word_embeddings)

    keywords_idx = [np.argmax(word_doc_similarity)]
    candidates_idx = [i for i in range(len(words)) if i != keywords_idx[0]]

    for _ in range(top_n - 1):
        candidate_similarities = word_doc_similarity[candidates_idx, :]
        target_similarities = np.max(word_similarity[candidates_idx][:, keywords_idx], axis=1)

        mmr = (lambda_param * candidate_similarities) - ((1-lambda_param) * target_similarities)
        mmr_idx = candidates_idx[np.argmax(mmr)]

        keywords_idx.append(mmr_idx)
        candidates_idx.remove(mmr_idx)

    return [words[idx] for idx in keywords_idx]

# Directory setup
base_directory = r"E:\DPI - Project\Questiones Bank\Data Analytics\Data Analyst Specialist\Data Analyst Specialist (revised)\(S4-S7) Make Data-Driven Decisions"

# Create necessary folders
create_folders(base_directory)

# Run the process and limit questions to 30
process_pptx_files(base_directory, base_directory, max_questions=30)
