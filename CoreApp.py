import os
import gc
import re
import torch
import random
import numpy as np
import pandas as pd
from pptx import Presentation
from sense2vec import Sense2Vec
from sentence_transformers import SentenceTransformer
import nltk
import string
import pke
import requests
from nltk.corpus import stopwords, wordnet as wn
from concurrent.futures import ThreadPoolExecutor

nltk.download('punkt')
nltk.download('wordnet')
nltk.download('stopwords')
nltk.download('omw-1.4')

# Load SpaCy model
import spacy
spacy_model = spacy.load('en_core_web_sm')

# Datamuse API URL
DATAMUSE_API_URL = "https://api.datamuse.com/words"

# Topic-specific fallback distractors (max two words each)
fallback_distractors = {
    "Software Development": [
        "Agile methods", 
        "Version control", 
        "Code review", 
        "Unit testing", 
        "SOLID principles", 
        "CI/CD pipeline", 
        "Code refactoring", 
        "Design patterns", 
        "Debugging tools", 
        "Static analysis"
    ],
    "Prompt Engineer": [
        "Context management", 
        "Model tuning", 
        "Token limit", 
        "Zero-shot", 
        "Few-shot", 
        "Prompt length", 
        "Output variance", 
        "Response generation", 
        "Word embeddings", 
        "Bias reduction"
    ],
    "Oracle ERP": [
        "General Ledger", 
        "Accounts Payable", 
        "Supply Chain", 
        "Project management", 
        "Financial reporting", 
        "Asset management", 
        "Procurement system", 
        "Order management", 
        "Inventory control", 
        "Expense tracking"
    ],
    "Odoo ERP": [
        "Sales orders", 
        "Customer CRM", 
        "Purchase orders", 
        "Inventory module", 
        "Accounting system", 
        "Timesheets tracking", 
        "Manufacturing orders", 
        "Human resources", 
        "Marketing automation", 
        "Payroll management"
    ],
    "Infrastructure and Security": [
        "Network firewall", 
        "Data encryption", 
        "Access control", 
        "Intrusion detection", 
        "Cloud security", 
        "Penetration testing", 
        "Backup recovery", 
        "Authentication methods", 
        "Threat monitoring", 
        "Malware protection"
    ],
    "Google Project Management": [
        "Project timeline", 
        "Risk analysis", 
        "Stakeholder communication", 
        "Team collaboration", 
        "Resource allocation", 
        "Budget tracking", 
        "Scope management", 
        "Agile workflow", 
        "Change management", 
        "Task prioritization"
    ],
    "Data Analytics": [
        "Data modeling", 
        "Predictive analytics", 
        "Descriptive statistics", 
        "Dashboard visualization", 
        "Data cleaning", 
        "Trend analysis", 
        "Data mining", 
        "Business insights", 
        "Ad-hoc reports", 
        "Big data"
    ],
    "AI & Data Science": [
        "Machine learning", 
        "Neural networks", 
        "Supervised learning", 
        "Unsupervised learning", 
        "Deep learning", 
        "Data preprocessing", 
        "Model evaluation", 
        "Natural language", 
        "Pattern recognition", 
        "Algorithm tuning"
    ],
    "Software Tester": [
        "Unit testing", 
        "Regression testing", 
        "Test automation", 
        "Integration tests", 
        "Load testing", 
        "Bug tracking", 
        "User acceptance", 
        "Manual testing", 
        "Functional testing", 
        "Test scripts"
    ],
    "React Web Developer": [
        "Component lifecycle", 
        "State management", 
        "Virtual DOM", 
        "Hooks API", 
        "Event handling", 
        "Functional components", 
        "React Router", 
        "JSX syntax", 
        "Context API", 
        "Redux store"
    ],
    "PHP Web Developer": [
        "Server scripting", 
        "MVC framework", 
        "Database queries", 
        "Form validation", 
        "Session management", 
        "API integration", 
        "Error handling", 
        "Authentication system", 
        "Cookie management", 
        "File upload"
    ],
    "Mobile App Developer": [
        "Cross-platform", 
        "UI components", 
        "App deployment", 
        "Push notifications", 
        "Touch interface", 
        "Mobile frameworks", 
        "API calls", 
        "Device storage", 
        "App debugging", 
        "User experience"
    ],
    "Frontend Developer": [
        "Responsive design", 
        "CSS grids", 
        "JavaScript frameworks", 
        "DOM manipulation", 
        "Web components", 
        "Flexbox layout", 
        "Cross-browser", 
        "HTML semantics", 
        "Web animations", 
        "SEO optimization"
    ],
    "DevOps Engineer": [
        "Continuous integration", 
        "Continuous deployment", 
        "Infrastructure code", 
        "Version control", 
        "Containerization tools", 
        "Server orchestration", 
        "Load balancing", 
        "Automated builds", 
        "Environment variables", 
        "Cloud deployments"
    ],
    ".Net Web Developer": [
        "ASP.NET Core", 
        "MVC pattern", 
        "Entity Framework", 
        "LINQ queries", 
        "Razor pages", 
        "Web API", 
        "Identity management", 
        "Dependency injection", 
        "Blazor components", 
        "C# programming"
    ],
    "PowerBI Engineer": [
        "Data models", 
        "Power Query", 
        "DAX formulas", 
        "Report dashboards", 
        "Data connections", 
        "Interactive visuals", 
        "SQL integration", 
        "Custom visuals", 
        "Real-time data", 
        "Workspace collaboration"
    ],
    "Data Analyst Specialist": [
        "Data cleaning", 
        "Business reporting", 
        "Statistical analysis", 
        "Trend forecasting", 
        "KPI tracking", 
        "Data warehousing", 
        "Ad-hoc reporting", 
        "Predictive models", 
        "SQL queries", 
        "Data mining"
    ],
    "Microsoft Machine Learning Engineer": [
        "Azure ML", 
        "Model deployment", 
        "Deep learning", 
        "Hyperparameter tuning", 
        "Model accuracy", 
        "Data pipelines", 
        "Neural networks", 
        "Data preprocessing", 
        "Model validation", 
        "Feature engineering"
    ],
    "Microsoft Data Engineer": [
        "ETL processes", 
        "Data lakes", 
        "Azure Data Factory", 
        "Big data", 
        "Data pipelines", 
        "Data warehouses", 
        "SQL queries", 
        "Data governance", 
        "Stream processing", 
        "Data security"
    ],
    "IBM Data Scientist": [
        "Data models", 
        "AI algorithms", 
        "Predictive analytics", 
        "Machine learning", 
        "Data visualization", 
        "Neural networks", 
        "Supervised learning", 
        "Unsupervised learning", 
        "Data preparation", 
        "Feature selection"
    ],
    "Generative AI": [
        "Text generation", 
        "Image synthesis", 
        "GAN models", 
        "Transformer networks", 
        "Style transfer", 
        "Conditional GAN", 
        "GPT models", 
        "Text embeddings", 
        "Content generation", 
        "AI creativity"
    ]
}

# Dynamic Question Templates for Various Tracks
question_templates = {
    # Software Development Track
    "Software Development": [
        "What is the purpose of {keyword} in software development?",
        "How does {keyword} improve the development process?",
        "Why is {keyword} essential for building scalable applications?",
        "What are the key principles of {keyword}?",
        "What challenges are addressed by {keyword}?",
        "How is {keyword} applied in Agile methodologies?",
        "What are the benefits of using {keyword} in software design?",
        "How does {keyword} ensure code maintainability?",
        "What impact does {keyword} have on software testing?",
        "How is {keyword} integrated into the CI/CD pipeline?"
    ],

    # Prompt Engineer Track
    "Prompt Engineer": [
        "What is the role of {keyword} in optimizing language models?",
        "How does {keyword} affect the output of a language model?",
        "Why is {keyword} crucial in prompt design?",
        "What techniques are used to improve {keyword} in prompt engineering?",
        "How can {keyword} influence the generation of model responses?",
        "What impact does {keyword} have on the efficiency of AI models?",
        "How does {keyword} improve the performance of language models?",
        "What are common mistakes made when designing prompts with {keyword}?",
        "Why is {keyword} a key factor in few-shot learning?",
        "What challenges arise when using {keyword} in GPT models?"
    ],

    # Oracle ERP Track
    "Oracle ERP": [
        "What are the main functionalities of {keyword} in Oracle ERP?",
        "How does {keyword} help in managing enterprise resources?",
        "Why is {keyword} an important feature of Oracle ERP systems?",
        "How does {keyword} improve business processes in Oracle ERP?",
        "What are the challenges of implementing {keyword} in Oracle ERP?",
        "How does {keyword} optimize supply chain management?",
        "What role does {keyword} play in financial management within Oracle ERP?",
        "How does {keyword} enhance project management capabilities?",
        "Why is {keyword} critical for HR operations in Oracle ERP?",
        "What are the benefits of using {keyword} in procurement processes?"
    ],

    # Odoo ERP Track
    "Odoo ERP": [
        "What are the core modules of {keyword} in Odoo ERP?",
        "How does {keyword} streamline operations in Odoo ERP?",
        "What benefits does {keyword} provide to small businesses using Odoo?",
        "How does {keyword} differ from other ERP systems in the market?",
        "What are the key challenges in implementing {keyword} within Odoo ERP?",
        "Why is {keyword} critical for managing inventory in Odoo ERP?",
        "How does {keyword} improve sales processes in Odoo?",
        "What role does {keyword} play in customer relationship management (CRM)?",
        "How is {keyword} used for financial reporting in Odoo?",
        "What are the customization options available for {keyword} in Odoo ERP?"
    ],

    # Infrastructure and Security Track
    "Infrastructure and Security": [
        "What role does {keyword} play in securing IT infrastructure?",
        "How does {keyword} improve data security in cloud environments?",
        "Why is {keyword} essential for protecting network traffic?",
        "What are the key components of {keyword} in a security architecture?",
        "How does {keyword} help prevent data breaches?",
        "What challenges are associated with implementing {keyword} in security?",
        "What impact does {keyword} have on disaster recovery planning?",
        "How does {keyword} ensure compliance with regulatory standards?",
        "What is the role of {keyword} in access control management?",
        "How does {keyword} detect and mitigate cyber threats?"
    ],

    # Google Project Management Track
    "Google Project Management": [
        "What is the significance of {keyword} in project planning?",
        "How does {keyword} help in managing project timelines?",
        "Why is {keyword} important for stakeholder communication?",
        "What are the best practices for using {keyword} in risk management?",
        "How does {keyword} improve project scope management?",
        "What challenges arise when using {keyword} in team collaboration?",
        "What are the key benefits of {keyword} for agile project management?",
        "How does {keyword} streamline resource allocation?",
        "What role does {keyword} play in project budgeting?",
        "How is {keyword} used for tracking project progress?"
    ],

    # Data Analytics Track
    "Data Analytics": [
        "What role does {keyword} play in data analytics?",
        "How does {keyword} improve data-driven decision-making?",
        "Why is {keyword} important for analyzing large datasets?",
        "What are the challenges of using {keyword} in data analysis?",
        "How does {keyword} enhance data visualization?",
        "What impact does {keyword} have on predictive analytics?",
        "How is {keyword} used in building dashboards for data insights?",
        "What are the limitations of {keyword} in business intelligence?",
        "Why is {keyword} important for real-time data analytics?",
        "How does {keyword} support data cleaning and preprocessing?"
    ],

    # AI & Data Science Track
    "AI & Data Science": [
        "What is the role of {keyword} in AI model development?",
        "How does {keyword} improve the accuracy of data science models?",
        "Why is {keyword} critical for deep learning applications?",
        "What challenges are associated with using {keyword} in machine learning?",
        "How does {keyword} support natural language processing tasks?",
        "What is the significance of {keyword} in AI research?",
        "How is {keyword} applied in neural networks?",
        "What are the benefits of using {keyword} for big data analysis?",
        "How does {keyword} contribute to AI model interpretability?",
        "What impact does {keyword} have on AI ethics and fairness?"
    ],

    # Software Tester Track
    "Software Tester": [
        "What role does {keyword} play in software testing?",
        "How does {keyword} improve the software testing process?",
        "Why is {keyword} important for ensuring software quality?",
        "What are the common challenges associated with {keyword} in testing?",
        "How does {keyword} help in identifying bugs in software?",
        "What impact does {keyword} have on test automation?",
        "Why is {keyword} essential for performance testing?",
        "How does {keyword} support user acceptance testing (UAT)?",
        "What are the benefits of using {keyword} in regression testing?",
        "What role does {keyword} play in integration testing?"
    ],

    # React Web Developer Track
    "React Web Developer": [
        "What is the purpose of {keyword} in React web development?",
        "How does {keyword} improve user interface design in React?",
        "Why is {keyword} important for managing state in React applications?",
        "What are the best practices for using {keyword} in React components?",
        "How does {keyword} support performance optimization in React apps?",
        "What role does {keyword} play in routing within a React application?",
        "What challenges arise when using {keyword} in React hooks?",
        "Why is {keyword} essential for React development with Redux?",
        "How does {keyword} contribute to React’s component-based architecture?",
        "What impact does {keyword} have on cross-browser compatibility in React?"
    ],

    # PHP Web Developer Track
    "PHP Web Developer": [
        "What is the purpose of {keyword} in PHP development?",
        "How does {keyword} improve server-side scripting in PHP?",
        "Why is {keyword} important for database management in PHP applications?",
        "What are the key features of {keyword} in modern PHP frameworks?",
        "How does {keyword} support API development in PHP?",
        "What challenges are associated with using {keyword} in PHP web security?",
        "What impact does {keyword} have on performance in PHP applications?",
        "How is {keyword} used for session management in PHP?",
        "Why is {keyword} critical for object-oriented programming in PHP?",
        "How does {keyword} help with error handling in PHP applications?"
    ],

    # Mobile App Developer Track
    "Mobile App Developer": [
        "What is the role of {keyword} in mobile app development?",
        "How does {keyword} enhance user experience in mobile apps?",
        "Why is {keyword} important for mobile app performance optimization?",
        "What are the key challenges in using {keyword} for mobile app development?",
        "How does {keyword} support cross-platform development?",
        "What impact does {keyword} have on mobile app security?",
        "How is {keyword} used for mobile app testing and debugging?",
        "Why is {keyword} essential for mobile app deployment?",
        "What are the benefits of using {keyword} for mobile app analytics?",
        "How does {keyword} improve mobile app monetization?"
    ],

    # Frontend Developer Track
    "Frontend Developer": [
        "What is the role of {keyword} in frontend development?",
        "How does {keyword} improve website performance?",
        "Why is {keyword} important for responsive web design?",
        "What challenges arise when using {keyword} in frontend frameworks?",
        "How does {keyword} contribute to better user interface design?",
        "Why is {keyword} critical for cross-browser compatibility?",
        "What are the benefits of using {keyword} for CSS animations?",
        "How does {keyword} enhance website accessibility?",
        "What impact does {keyword} have on SEO in frontend development?",
        "What role does {keyword} play in frontend build tools?"
    ],

    # DevOps Engineer Track
    "DevOps Engineer": [
        "What is the purpose of {keyword} in DevOps?",
        "How does {keyword} streamline continuous integration and deployment?",
        "Why is {keyword} essential for automation in DevOps pipelines?",
        "What challenges arise when using {keyword} in DevOps workflows?",
        "How does {keyword} improve collaboration between development and operations?",
        "What role does {keyword} play in monitoring and logging in DevOps?",
        "How does {keyword} contribute to infrastructure as code?",
        "Why is {keyword} critical for containerization in DevOps?",
        "How does {keyword} support version control in a DevOps environment?",
        "What are the benefits of using {keyword} for cloud-native DevOps?"
    ],

    # .Net Web Developer Track
    ".Net Web Developer": [
        "What is the role of {keyword} in .NET web development?",
        "How does {keyword} improve the performance of .NET applications?",
        "Why is {keyword} important for web services in .NET?",
        "What are the key features of {keyword} in ASP.NET Core?",
        "How does {keyword} support API development in .NET?",
        "What challenges arise when using {keyword} in .NET security?",
        "Why is {keyword} essential for dependency injection in .NET?",
        "How does {keyword} help with unit testing in .NET applications?",
        "What role does {keyword} play in Entity Framework for .NET?",
        "How does {keyword} enhance cross-platform development in .NET?"
    ],

    # PowerBI Engineer Track
    "PowerBI Engineer": [
        "What is the role of {keyword} in PowerBI reports?",
        "How does {keyword} improve data visualization in PowerBI?",
        "Why is {keyword} critical for PowerBI data modeling?",
        "What are the best practices for using {keyword} in PowerBI dashboards?",
        "How does {keyword} enhance PowerBI’s performance for large datasets?",
        "What role does {keyword} play in PowerBI’s report publishing?",
        "Why is {keyword} important for PowerBI’s integration with other data sources?",
        "How does {keyword} support PowerBI’s collaboration features?",
        "What challenges are associated with using {keyword} in PowerBI’s DAX formulas?",
        "How does {keyword} help in creating PowerBI’s real-time dashboards?"
    ],

    # Data Analyst Specialist Track
    "Data Analyst Specialist": [
        "What is the role of {keyword} in data analysis?",
        "How does {keyword} improve decision-making processes?",
        "Why is {keyword} important for cleaning and preprocessing data?",
        "What are the common challenges of using {keyword} in data analysis?",
        "How does {keyword} enhance the quality of data visualizations?",
        "What impact does {keyword} have on predictive analytics models?",
        "Why is {keyword} critical for building business intelligence reports?",
        "How does {keyword} contribute to understanding data trends?",
        "What role does {keyword} play in summarizing large datasets?",
        "What are the benefits of using {keyword} for ad-hoc reporting?"
    ],

    # Microsoft Machine Learning Engineer Track
    "Microsoft Machine Learning Engineer": [
        "What is the purpose of {keyword} in machine learning models?",
        "How does {keyword} improve the accuracy of predictions?",
        "Why is {keyword} important for model training on large datasets?",
        "What challenges are associated with using {keyword} in Azure ML?",
        "How does {keyword} enhance the performance of neural networks?",
        "What role does {keyword} play in hyperparameter tuning for ML models?",
        "Why is {keyword} critical for deploying models in the cloud?",
        "How does {keyword} help in monitoring deployed ML models?",
        "What impact does {keyword} have on the scalability of machine learning solutions?",
        "How does {keyword} support model evaluation and validation?"
    ],

    # Microsoft Data Engineer Track
    "Microsoft Data Engineer": [
        "What is the role of {keyword} in data engineering?",
        "How does {keyword} improve data pipeline automation?",
        "Why is {keyword} important for data storage and management?",
        "What are the challenges of using {keyword} in distributed data systems?",
        "How does {keyword} enhance the performance of ETL processes?",
        "What role does {keyword} play in data integration and transformation?",
        "Why is {keyword} critical for building scalable data infrastructures?",
        "How does {keyword} support data governance and compliance?",
        "What impact does {keyword} have on big data analytics?",
        "How is {keyword} used to ensure data quality and integrity?"
    ],

    # IBM Data Scientist Track
    "IBM Data Scientist": [
        "What is the purpose of {keyword} in data science?",
        "How does {keyword} improve the accuracy of data models?",
        "Why is {keyword} important for data exploration and analysis?",
        "What challenges arise when using {keyword} in big data projects?",
        "How does {keyword} support data-driven decision-making?",
        "What role does {keyword} play in predictive modeling?",
        "Why is {keyword} critical for building machine learning models?",
        "How does {keyword} contribute to the understanding of data trends?",
        "What impact does {keyword} have on AI-powered data analysis?",
        "How is {keyword} used to evaluate the performance of data science models?"
    ],

    # Generative AI Track
    "Generative AI": [
        "What is the role of {keyword} in generative AI models?",
        "How does {keyword} improve the quality of generated content?",
        "Why is {keyword} important for training generative models?",
        "What challenges arise when using {keyword} in GANs?",
        "How does {keyword} support text generation in GPT models?",
        "What are the key techniques used in {keyword} for content creation?",
        "Why is {keyword} critical for generating realistic images?",
        "How does {keyword} help in the creative process of generative AI?",
        "What role does {keyword} play in improving generative adversarial networks?",
        "How does {keyword} enhance the diversity of generated outputs?"
    ]
}


# Initialize models
def initialize_models():
    device = torch.device("cuda" if torch.cuda.is_available() else "cpu")
    print(f"Using device: {device}")
    
    s2v = Sense2Vec().from_disk(r'C:\Python\Lib\site-packages\sense2vec\tests\data')
    sentence_transformer_model = SentenceTransformer('sentence-transformers/msmarco-distilbert-base-v2')
    
    return s2v, sentence_transformer_model, device

def create_folders(base_directory):
    for folder in ["Quizzes", "Question Bank", "Weekly Assignments"]:
        os.makedirs(os.path.join(base_directory, folder), exist_ok=True)

def read_pptx(file_path):
    presentation = Presentation(file_path)
    return [
        "\n".join(shape.text for shape in slide.shapes if shape.has_text_frame)
        for slide in presentation.slides
    ]

def clean_text(text):
    return re.sub(r'\s+', ' ', re.sub(r'[^\x09\x0A\x0D\x20-\x7F]', '', text)).strip()

def get_keywords(content, domain_stopwords=[]):
    extractor = pke.unsupervised.MultipartiteRank()
    extractor.load_document(input=content, language='en', spacy_model=spacy_model)
    stoplist = list(string.punctuation) + stopwords.words('english') + domain_stopwords
    extractor.candidate_selection(pos={'NOUN', 'PROPN'})
    extractor.candidate_weighting(alpha=1.1, threshold=0.75, method='average')
    keyphrases = extractor.get_n_best(n=10)
    
    return [val[0] for val in keyphrases if len(val[0].split()) >= 2 and len(val[0]) > 2 and not any(char.isdigit() for char in val[0])]

def get_datamuse_related_words(word, max_results=5):
    response = requests.get(DATAMUSE_API_URL, params={'ml': word, 'max': max_results})
    return [item['word'] for item in response.json()] if response.status_code == 200 else []

def sense2vec_get_words(word, s2v, top_n=10):
    word_with_pos = f"{word.lower()}|NOUN"
    return [w[0].split("|")[0] for w in s2v.most_similar(word_with_pos, n=top_n)] if word_with_pos in s2v else []

def get_wordnet_distractors(word, pos=wn.NOUN):
    return list(set(
        lemma.name().replace('_', ' ')
        for syn in wn.synsets(word, pos=pos)
        for lemma in syn.lemmas()
        if lemma.name().lower() != word.lower()
    ))[:5]

def get_distractors(word, topic, s2v, sentence_transformer_model, top_n=40):
    # Attempt to gather distractors from different sources (Sense2Vec, Datamuse, WordNet)
    distractors = list(set(
        sense2vec_get_words(word, s2v, top_n) +
        get_datamuse_related_words(word) +
        get_wordnet_distractors(word)
    ))
    
    # Filter out the word itself from the distractors
    distractors = [d for d in distractors if d.lower() != word.lower()]

    # If fewer than 3 distractors, use fallback distractors specific to the topic
    if len(distractors) < 3:
        distractors = random.sample(fallback_distractors[topic], 3)
    
    return distractors[:3]

def generate_dynamic_question_templates(keyword, topic, context=None):
    related_term = random.choice(get_datamuse_related_words(keyword) or ['[related term]'])
    
    if topic in question_templates:
        return random.choice(question_templates[topic]).format(keyword=keyword)
    else:
        return f"What is the role of {keyword} in {topic}?"

def process_slide(slide_text, s2v, sentence_transformer_model, question_count, max_questions=20, topic=None):
    cleaned_text = clean_text(slide_text)
    imp_keywords = get_keywords(cleaned_text)
    
    slide_questions = []
    for answer in imp_keywords:
        if question_count[0] >= max_questions:
            break
        question = generate_dynamic_question_templates(answer, topic)
        distractors = get_distractors(answer.capitalize(), question, s2v, sentence_transformer_model)
        if not distractors:
            continue

        options = random.sample(distractors + [answer], 4)
        correct_option = options.index(answer)
        slide_questions.append({
            "Question": question,
            "Option A": options[0],
            "Option B": options[1],
            "Option C": options[2],
            "Option D": options[3],
            "Answer": chr(65 + correct_option)
        })
        question_count[0] += 1

    return slide_questions

def generate_questions_for_pptx(pptx_path, output_directory, s2v, sentence_transformer_model, max_questions=20, topic=None):
    all_questions = []
    question_count = [0]
    slide_texts = read_pptx(pptx_path)

    with ThreadPoolExecutor(max_workers=4) as executor:
        results = executor.map(
            lambda slide_text: process_slide(slide_text, s2v, sentence_transformer_model, question_count, max_questions, topic),
            slide_texts
        )
        for slide_questions in results:
            all_questions.extend(slide_questions)
            if question_count[0] >= max_questions:
                break

    if all_questions:
        save_to_excel(all_questions, pptx_path, output_directory)

def save_to_excel(questions, pptx_path, output_directory):
    df = pd.DataFrame(questions)
    
    # Drop duplicates based on the "Question" column
    df_cleaned = df.drop_duplicates(subset=['Question'])

    excel_filename = os.path.basename(pptx_path).replace('.pptx', '.xlsx')
    excel_path = os.path.join(output_directory, "Quizzes", excel_filename)
    df.to_excel(excel_path, index=False)
    print(f"Processed {pptx_path} -> {excel_path}")

def merge_all_quizzes(output_directory):
    quizzes_folder = os.path.join(output_directory, "Quizzes")
    question_bank_folder = os.path.join(output_directory, "Question Bank")
    
    all_files = [os.path.join(quizzes_folder, f) for f in os.listdir(quizzes_folder) if f.endswith('.xlsx')]
    
    merged_df = pd.concat([pd.read_excel(f) for f in all_files], ignore_index=True)
    
    merged_file_path = os.path.join(question_bank_folder, 'Merged_Quiz_Question_Bank.xlsx')
    merged_df.to_excel(merged_file_path, index=False)
    
    print(f"Merged quiz saved at: {merged_file_path}")
    return merged_file_path

def process_pptx_files(directory, output_directory, max_questions=20, topic=None):
    s2v, sentence_transformer_model, _ = initialize_models()

    pptx_files = [
        os.path.join(root, file)
        for root, _, files in os.walk(directory)
        for file in files if file.endswith(".pptx")
    ]

    create_folders(output_directory)

    with ThreadPoolExecutor(max_workers=2) as executor:
        executor.map(
            lambda pptx_path: generate_questions_for_pptx(pptx_path, output_directory, s2v, sentence_transformer_model, max_questions, topic),
            pptx_files
        )

    merge_all_quizzes(output_directory)
    gc.collect()

if __name__ == "__main__":
    base_directory = r"E:\DPI - Project\Questiones Bank\Data Analytics\PowerBI Engineer\PowerBI Engineer (Revised)\(S1-S4) Introduction to Data Analysis"
    topic = "PowerBI Engineer"  # Set the appropriate topic
    process_pptx_files(base_directory, base_directory, max_questions=20, topic=topic)
