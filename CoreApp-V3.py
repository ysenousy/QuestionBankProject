import os
import gc
import re
import torch
import random
import numpy as np
from pptx import Presentation
from docx import Document
from sense2vec import Sense2Vec
from sentence_transformers import SentenceTransformer, util
import nltk
import string
import pke
from nltk.corpus import stopwords
from concurrent.futures import ThreadPoolExecutor

nltk.download('punkt')
nltk.download('wordnet')
nltk.download('stopwords')
nltk.download('omw-1.4')

# Load SpaCy model
import spacy
spacy_model = spacy.load('en_core_web_sm')

# Initialize models
def initialize_models():
    device = torch.device("cuda" if torch.cuda.is_available() else "cpu")
    print(f"Using device: {device}")

    # Load Sense2Vec model
    s2v = Sense2Vec().from_disk(r'C:\Python\Lib\site-packages\sense2vec\tests\data')

    # Load Sentence Transformer model
    sentence_transformer_model = SentenceTransformer('sentence-transformers/msmarco-distilbert-base-v2')

    return s2v, sentence_transformer_model, device

# Create necessary folders
def create_folders(base_directory):
    os.makedirs(os.path.join(base_directory, "Quizzes"), exist_ok=True)
    os.makedirs(os.path.join(base_directory, "Question Bank"), exist_ok=True)
    os.makedirs(os.path.join(base_directory, "Weekly Assignments"), exist_ok=True)

# Read all text from a PowerPoint file
def read_pptx(file_path):
    presentation = Presentation(file_path)
    all_text = []

    for slide in presentation.slides:
        slide_text = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                slide_text.append(paragraph.text)
        all_text.append("\n".join(slide_text))

    return all_text

# Clean input text
def clean_text(text):
    text = re.sub(r'[^\x09\x0A\x0D\x20-\x7F]', '', text)
    text = re.sub(r'\d{1,2}/\d{1,2}/\d{2,4}', '', text)
    text = re.sub(r'\s+', ' ', text)
    return text.strip()

# Enhanced keyword extraction
def get_keywords(content):
    extractor = pke.unsupervised.MultipartiteRank()
    extractor.load_document(input=content, language='en', spacy_model=spacy_model)
    pos = {'PROPN', 'NOUN', 'ADJ', 'VERB'}  # Focus on nouns and verbs
    stoplist = list(string.punctuation) + stopwords.words('english')
    extractor.candidate_selection(pos=pos)
    extractor.candidate_weighting(alpha=1.1, threshold=0.75, method='average')
    keyphrases = extractor.get_n_best(n=10)  # Limit to 10 top phrases
    return [val[0] for val in keyphrases if len(val[0].split()) > 1]  # Avoid single-word keywords

# Generate distractors using Sense2Vec
def sense2vec_get_words(word, s2v, top_n=10, sentence=None):
    word = word.lower()
    word_with_pos = f"{word}|NOUN"
    if word_with_pos in s2v:
        most_similar = s2v.most_similar(word_with_pos, n=top_n)
        return [w[0].split("|")[0] for w in most_similar]
    else:
        return []

# Function to filter semantically similar distractors
def filter_semantically_relevant_distractors(distractors, context, model, threshold=0.7):
    context_embedding = model.encode(context, convert_to_tensor=True)
    relevant_distractors = []
    for distractor in distractors:
        distractor_embedding = model.encode(distractor, convert_to_tensor=True)
        if util.pytorch_cos_sim(context_embedding, distractor_embedding) > threshold:
            relevant_distractors.append(distractor)
    return relevant_distractors

# Generate meaningful distractors using Sense2Vec
def get_distractors(word, sentence, s2v, sentence_transformer_model, top_n=40, lambda_val=0.2):
    original_word = word.lower()
    distractors = sense2vec_get_words(word, s2v, top_n, sentence)
    # Remove the correct answer and duplicates
    filtered_distractors = list(set([distractor for distractor in distractors if distractor.lower() != original_word]))
    if len(filtered_distractors) < 3:
        return []  # Not enough distractors, consider returning none or fallback to another method
    distractors_new = filtered_distractors

    keyword_embedding = sentence_transformer_model.encode([sentence + " " + word.capitalize()])
    distractor_embeddings = sentence_transformer_model.encode(distractors_new)

    max_keywords = min(len(distractors_new), 5)
    filtered_keywords = mmr(keyword_embedding, distractor_embeddings, distractors_new, max_keywords, lambda_val)
    
    return [w.capitalize() for w in filtered_keywords if w.lower() != original_word]

# Enhanced question templates
def generate_question_templates(keyword):
    templates = [
        f"What is the primary function of {keyword}?",
        f"How does {keyword} contribute to the overall process?",
        f"In what way does {keyword} impact the system?",
        f"What role does {keyword} play in achieving the desired outcome?",
        f"Why is {keyword} considered important in this context?",
        f"What are the key characteristics of {keyword}?",
        f"How would you describe the relationship between {keyword} and other related concepts?",
        f"What are the potential consequences of neglecting {keyword}?",
        f"In what scenarios is {keyword} most relevant?",
        f"How has the understanding of {keyword} evolved over time?",
        f"What are some common misconceptions about {keyword}?",
        f"How can {keyword} be effectively implemented or utilized?",
        f"What are the limitations or challenges associated with {keyword}?",
        f"How does {keyword} compare to alternative approaches or concepts?",
        f"What future developments or trends are expected in relation to {keyword}?"
    ]
    return random.choice(templates)

# Process a single slide (parallel processing)
def process_slide(slide_text, s2v, sentence_transformer_model, device, question_count, max_questions=20):
    cleaned_text = clean_text(slide_text)
    imp_keywords = get_keywords(cleaned_text)
    
    slide_questions = []
    for answer in imp_keywords:
        if question_count[0] >= max_questions:  # Stop when 20 questions are reached
            break
        question = generate_question_templates(answer)
        distractors = get_distractors(answer.capitalize(), question, s2v, sentence_transformer_model)
        if len(distractors) == 0:
            distractors = imp_keywords

        random.shuffle(distractors)
        options = distractors[:3] + [answer]
        random.shuffle(options)
        correct_option = options.index(answer)
        output = f"Question: {question}\n"
        for i, opt in enumerate(options):
            output += f"{chr(65 + i)}) {opt}\n"
        output += f"Answer: {chr(65 + correct_option)}\n\n"
        slide_questions.append(output)
        question_count[0] += 1  # Increment the global question counter

    return slide_questions

# Generate all questions for a PPTX
def generate_questions_for_pptx(pptx_path, output_directory, s2v, sentence_transformer_model, device, max_questions=20):
    all_questions = []
    question_count = [0]  # Use list to allow modification inside process_slide
    slide_texts = read_pptx(pptx_path)

    # Parallel processing of slides using ThreadPoolExecutor
    with ThreadPoolExecutor(max_workers=4) as executor:
        results = executor.map(lambda slide_text: process_slide(slide_text, s2v, sentence_transformer_model, device, question_count, max_questions), slide_texts)
        for slide_questions in results:
            all_questions.extend(slide_questions)
            if question_count[0] >= max_questions:  # Stop once we have enough questions
                break

    full_text = "\n\n".join(all_questions)
    word_filename = os.path.basename(pptx_path).replace('.pptx', '.docx')
    word_path = os.path.join(output_directory, "Quizzes", word_filename)
    create_word_document(full_text, word_path)

    print(f"Processed {pptx_path} -> {word_path}")

# Process PowerPoint files and generate quizzes
def process_pptx_files(directory, output_directory, max_questions=20):
    s2v, sentence_transformer_model, device = initialize_models()

    pptx_files = [os.path.join(root, file) 
                  for root, _, files in os.walk(directory) 
                  for file in files if file.endswith(".pptx")]

    # Use ThreadPoolExecutor to process each PPTX file in parallel
    with ThreadPoolExecutor(max_workers=2) as executor:
        executor.map(lambda pptx_path: generate_questions_for_pptx(pptx_path, output_directory, s2v, sentence_transformer_model, device, max_questions), pptx_files)

    del s2v, sentence_transformer_model
    gc.collect()

# Utility functions for formatting text and saving documents
def create_word_document(questions, output_path):
    doc = Document()
    doc.add_heading('Generated Questions', 0)
    doc.add_paragraph(questions)
    doc.save(output_path)

# Directory setup
base_directory = r"E:\DPI - Project\Questiones Bank\Data Analytics\Data Analyst Specialist\Data Analyst Specialist (revised)\(S4-S7) Make Data-Driven Decisions"

# Create necessary folders
create_folders(base_directory)

# Run the process and limit questions to 20
process_pptx_files(base_directory, base_directory, max_questions=20)
