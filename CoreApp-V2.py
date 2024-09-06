import os
import gc
import re
import torch
import random
import numpy as np
from pptx import Presentation
from docx import Document
from transformers import T5ForConditionalGeneration, T5Tokenizer
from sense2vec import Sense2Vec
from sentence_transformers import SentenceTransformer
from sklearn.metrics.pairwise import cosine_similarity
import nltk
from nltk.corpus import wordnet as wn
from nltk.tokenize import sent_tokenize
import string
import pke
from flashtext import KeywordProcessor
from collections import OrderedDict
from strsimpy.normalized_levenshtein import NormalizedLevenshtein
import spacy
from concurrent.futures import ThreadPoolExecutor

nltk.download('punkt')
nltk.download('brown')
nltk.download('wordnet')
nltk.download('stopwords')
nltk.download('omw-1.4')
from nltk.corpus import stopwords

# Load SpaCy model
spacy_model = spacy.load('en_core_web_sm')

# Initialize models and tokenizers
def initialize_models():
    device = torch.device("cuda" if torch.cuda.is_available() else "cpu")
    print(f"Using device: {device}")

    # Load Sense2Vec model
    s2v = Sense2Vec().from_disk(r'C:\Python\Lib\site-packages\sense2vec\tests\data')

    # Load T5 summarization model and tokenizer with explicit max length
    summary_model = T5ForConditionalGeneration.from_pretrained('t5-base').to(device)
    summary_tokenizer = T5Tokenizer.from_pretrained('t5-base', model_max_length=512)

    # Load T5 question generation model and tokenizer with explicit max length
    question_model = T5ForConditionalGeneration.from_pretrained('ramsrigouthamg/t5_squad_v1').to(device)
    question_tokenizer = T5Tokenizer.from_pretrained('ramsrigouthamg/t5_squad_v1', model_max_length=512)

    # Load Sentence Transformer model
    sentence_transformer_model = SentenceTransformer('sentence-transformers/msmarco-distilbert-base-v2')

    return s2v, summary_model, summary_tokenizer, question_model, question_tokenizer, sentence_transformer_model, device

# Function to create necessary folders
def create_folders(base_directory):
    os.makedirs(os.path.join(base_directory, "Quizzes"), exist_ok=True)
    os.makedirs(os.path.join(base_directory, "Question Bank"), exist_ok=True)
    os.makedirs(os.path.join(base_directory, "Weekly Assignments"), exist_ok=True)

# Function to read all text from a PowerPoint file
def read_pptx(file_path):
    presentation = Presentation(file_path)
    all_text = []

    for slide in presentation.slides:
        slide_text = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                slide_text.append(paragraph.text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            slide_text.append(paragraph.text)
        all_text.append("\n".join(slide_text))

    return all_text

# Function to clean input text thoroughly
def clean_text(text):
    text = re.sub(r'[^\x09\x0A\x0D\x20-\x7F]', '', text)
    text = re.sub(r'\d{1,2}/\d{1,2}/\d{2,4}', '', text)
    text = re.sub(r'\s+', ' ', text)
    return text.strip()

# Function to summarize text
def summarizer(text, model, tokenizer, device):
    text = text.strip().replace("\n", " ")
    text = "summarize: " + text
    max_len = 512
    encoding = tokenizer.encode_plus(text, max_length=max_len, pad_to_max_length=False, truncation=True, return_tensors="pt").to(device)
    input_ids, attention_mask = encoding["input_ids"], encoding["attention_mask"]

    outs = model.generate(input_ids=input_ids, attention_mask=attention_mask, early_stopping=True, num_beams=3, num_return_sequences=1, no_repeat_ngram_size=2, min_length=75, max_length=300)

    summary = tokenizer.decode(outs[0], skip_special_tokens=True)
    summary = postprocesstext(summary)
    return summary.strip()

# Function to extract keywords using pke
def get_keywords(content):
    extractor = pke.unsupervised.MultipartiteRank()
    extractor.load_document(input=content, language='en', spacy_model=spacy_model)
    pos = {'PROPN', 'NOUN', 'ADJ', 'VERB', 'ADP', 'ADV', 'DET', 'CONJ', 'NUM', 'PRON', 'X'}
    stoplist = list(string.punctuation)
    stoplist += stopwords.words('english')
    extractor.candidate_selection(pos=pos)
    extractor.candidate_weighting(alpha=1.1, threshold=0.75, method='average')
    keyphrases = extractor.get_n_best(n=15)
    return [val[0] for val in keyphrases]

# Function to generate questions using T5 model
def get_question(context, answer, model, tokenizer, device):
    text = f"context: {context} answer: {answer}"
    encoding = tokenizer.encode_plus(text, max_length=384, pad_to_max_length=False, truncation=True, return_tensors="pt").to(device)
    input_ids, attention_mask = encoding["input_ids"], encoding["attention_mask"]

    outs = model.generate(input_ids=input_ids, attention_mask=attention_mask, early_stopping=True, num_beams=5, num_return_sequences=1, no_repeat_ngram_size=2, max_length=72)

    question = tokenizer.decode(outs[0], skip_special_tokens=True)
    return question.replace("question:", "").strip()

# Function to retrieve similar words using Sense2Vec
def sense2vec_get_words(word, s2v, top_n=10, sentence=None):
    word = word.lower()
    word_with_pos = f"{word}|NOUN"  # Assuming the word is a noun; change as needed
    if word_with_pos in s2v:
        most_similar = s2v.most_similar(word_with_pos, n=top_n)
        return [w[0].split("|")[0] for w in most_similar]
    else:
        return []

# Function to generate distractors using Sense2Vec
def get_distractors(word, sentence, s2v, sentence_transformer_model, top_n=40, lambda_val=0.2):
    distractors = sense2vec_get_words(word, s2v, top_n, sentence)
    if len(distractors) == 0:
        return distractors
    distractors_new = [word.capitalize()]
    distractors_new.extend(distractors)

    keyword_embedding = sentence_transformer_model.encode([sentence + " " + word.capitalize()])
    distractor_embeddings = sentence_transformer_model.encode(distractors_new)

    max_keywords = min(len(distractors_new), 5)
    filtered_keywords = mmr(keyword_embedding, distractor_embeddings, distractors_new, max_keywords, lambda_val)
    
    return [w.capitalize() for w in filtered_keywords if w.lower() != word.lower()]

# Function to process a single slide (used for parallel processing)
def process_slide(slide_text, summary_model, summary_tokenizer, question_model, question_tokenizer, s2v, sentence_transformer_model, device):
    cleaned_text = clean_text(slide_text)
    summarized_text = summarizer(cleaned_text, summary_model, summary_tokenizer, device)
    imp_keywords = get_keywords(summarized_text)
    
    slide_questions = []
    for answer in imp_keywords:
        question = get_question(summarized_text, answer, question_model, question_tokenizer, device)
        distractors = get_distractors(answer.capitalize(), question, s2v, sentence_transformer_model)
        if len(distractors) == 0:
            distractors = imp_keywords

        random.shuffle(distractors)
        options = distractors[:3] + [answer]
        random.shuffle(options)
        correct_option = options.index(answer)
        output = f"Question: {question}\n"
        for i, opt in enumerate(options):
            output += f"{chr(65 + i)}) {opt}\n"
        output += f"Answer: {chr(65 + correct_option)}\n\n"
        slide_questions.append(output)

    return slide_questions

# Function to generate all questions for a PPTX
def generate_questions_for_pptx(pptx_path, output_directory, s2v, summary_model, summary_tokenizer, question_model, question_tokenizer, sentence_transformer_model, device):
    all_questions = []
    slide_texts = read_pptx(pptx_path)

    # Parallel processing of slides using ThreadPoolExecutor
    with ThreadPoolExecutor(max_workers=4) as executor:
        results = executor.map(lambda slide_text: process_slide(slide_text, summary_model, summary_tokenizer, question_model, question_tokenizer, s2v, sentence_transformer_model, device), slide_texts)
        for slide_questions in results:
            all_questions.extend(slide_questions)

    full_text = "\n\n".join(all_questions)
    word_filename = os.path.basename(pptx_path).replace('.pptx', '.docx')
    word_path = os.path.join(output_directory, "Quizzes", word_filename)
    create_word_document(full_text, word_path)

    print(f"Processed {pptx_path} -> {word_path}")

# Main function to process PowerPoint files and generate quizzes
def process_pptx_files(directory, output_directory):
    s2v, summary_model, summary_tokenizer, question_model, question_tokenizer, sentence_transformer_model, device = initialize_models()

    pptx_files = [os.path.join(root, file) 
                  for root, _, files in os.walk(directory) 
                  for file in files if file.endswith(".pptx")]

    # Use ThreadPoolExecutor to process each PPTX file in parallel
    with ThreadPoolExecutor(max_workers=2) as executor:
        executor.map(lambda pptx_path: generate_questions_for_pptx(pptx_path, output_directory, s2v, summary_model, summary_tokenizer, question_model, question_tokenizer, sentence_transformer_model, device), pptx_files)

    del s2v, summary_model, summary_tokenizer, question_model, question_tokenizer, sentence_transformer_model
    gc.collect()

# Utility functions for post-processing and formatting text
def postprocesstext(content):
    final = ""
    for sent in sent_tokenize(content):
        sent = sent.capitalize()
        final += " " + sent
    return final.strip()

def create_word_document(questions, output_path):
    doc = Document()
    doc.add_heading('Generated Questions', 0)
    doc.add_paragraph(questions)
    doc.save(output_path)

# Directory setup
base_directory = r"E:\DPI - Project\Questiones Bank\Data Analytics\Data Analyst Specialist\Data Analyst Specialist (revised)\(S4-S7) Make Data-Driven Decisions"

# Create necessary folders
create_folders(base_directory)

# Run the process
process_pptx_files(base_directory, base_directory)
